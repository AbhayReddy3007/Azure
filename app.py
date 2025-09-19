# app.py -- single-file Streamlit app with feedback editing for outlines
import os, re, tempfile, copy
import fitz, docx
import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from PIL import Image
import vertexai
from vertexai.generative_models import GenerativeModel
from google.oauth2 import service_account

# ---------------- GCP AUTH ----------------
# Make sure you have .streamlit/secrets.toml with [gcp_service_account] as discussed
gcp_service_account = st.secrets["gcp_service_account"]
credentials = service_account.Credentials.from_service_account_info(dict(gcp_service_account))
PROJECT_ID = gcp_service_account["project_id"]
REGION = "us-central1"
vertexai.init(project=PROJECT_ID, location=REGION, credentials=credentials)

# ---------------- MODEL ----------------
TEXT_MODEL_NAME = "gemini-2.0-flash"
TEXT_MODEL = GenerativeModel(TEXT_MODEL_NAME)

# ---------------- HELPERS / LLM CALLS ----------------
def call_vertex(prompt: str) -> str:
    try:
        response = TEXT_MODEL.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        # Provide a readable error for the UI
        return f"⚠️ Vertex AI error: {e}"

def extract_text(path: str, filename: str) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        text_parts = []
        doc = fitz.open(path)
        try:
            for page in doc:
                text_parts.append(page.get_text("text"))
        finally:
            doc.close()
        return "\n".join(text_parts)
    if name.endswith(".docx"):
        d = docx.Document(path)
        return "\n".join(p.text for p in d.paragraphs)
    if name.endswith(".txt"):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    return ""

def split_text(text: str, chunk_size: int = 8000, overlap: int = 300):
    if not text:
        return []
    chunks, start, n = [], 0, len(text)
    while start < n:
        end = min(start + chunk_size, n)
        chunks.append(text[start:end])
        if end == n:
            break
        start = max(0, end - overlap)
    return chunks

def summarize_long_text(full_text: str) -> str:
    chunks = split_text(full_text)
    if len(chunks) <= 1:
        return call_vertex(f"Summarize the following text in detail:\n\n{full_text}")
    partial_summaries = []
    for idx, ch in enumerate(chunks, start=1):
        mapped = call_vertex(f"Summarize this part of a longer document:\n\n{ch}")
        partial_summaries.append(f"Chunk {idx}:\n{mapped.strip()}")
    combined = "\n\n".join(partial_summaries)
    return call_vertex(f"Combine these summaries into one clean, well-structured summary:\n\n{combined}")

def generate_title(summary: str) -> str:
    prompt = f"""Read the following summary and create a short, clear, presentation-style title.
- Keep it under 10 words
- Do not include birth dates, long sentences, or excessive details
- Just give a clean title, like a presentation heading

Summary:
{summary}
"""
    return call_vertex(prompt).strip()

def parse_points(points_text: str):
    points = []
    current_title, current_content = None, []
    lines = [re.sub(r"[#*>`]", "", ln).rstrip() for ln in points_text.splitlines()]

    for line in lines:
        if not line or "Would you like" in line:
            continue
        m = re.match(r"^\s*(Slide|Section)\s*(\d+)\s*:\s*(.+)$", line, re.IGNORECASE)
        if m:
            if current_title:
                points.append({"title": current_title, "description": "\n".join(current_content)})
            current_title, current_content = m.group(3).strip(), []
            continue
        if line.strip().startswith("-"):
            text = line.lstrip("-").strip()
            if text:
                current_content.append(f"• {text}")
        elif line.strip().startswith(("•", "*")) or line.startswith("  "):
            text = line.lstrip("•*").strip()
            if text:
                current_content.append(f"- {text}")
        else:
            if line.strip():
                current_content.append(line.strip())

    if current_title:
        points.append({"title": current_title, "description": "\n".join(current_content)})
    return points

def generate_outline_from_desc(description: str, num_items: int = None):
    # If num_items is None, LLM decides slide count.
    if num_items:
        prompt = f"""Create a PowerPoint outline on: {description}.
Generate exactly {num_items} content slides (excluding the title slide).
Do NOT include a title slide — I will handle it separately.
Format strictly like this:
Slide 1: <Title>
- Bullet
- Bullet
- Bullet
"""
    else:
        prompt = f"""Create a PowerPoint outline on: {description}.
Decide the most appropriate number of content slides (excluding the title slide).
Each slide should have a short title and 3–4 bullet points.
Do NOT include a title slide — I will handle it separately.
Format strictly like this:
Slide 1: <Title>
- Bullet
- Bullet
- Bullet
"""
    points_text = call_vertex(prompt)
    return parse_points(points_text)

def edit_outline_with_feedback(outline: dict, feedback: str):
    """
    Uses the LLM to apply feedback and return an improved outline dict with 'title' and 'slides' keys.
    """
    outline_text = "\n".join([f"Slide {i+1}: {s['title']}\n{s['description']}" for i, s in enumerate(outline.get("slides", []))])
    prompt = f"""
You are an assistant improving a PowerPoint outline.

Current Outline:
Title: {outline.get('title', '')}
{outline_text}

Feedback:
{feedback}

Task:
- Apply the feedback to refine/improve the outline.
- Return the updated outline with the same format:
  Slide 1: <Title>
  - Bullet
  - Bullet
- Do NOT add a title slide.
"""
    updated = call_vertex(prompt)
    slides = parse_points(updated)
    # Keep the title (allow the user to change separately in UI), but return slides updated
    return {"title": outline.get("title", ""), "slides": slides}

# ---------------- PPT GENERATOR ----------------
def clean_title_text(title: str) -> str:
    return re.sub(r"\s+", " ", title.strip()) if title else "Presentation"

def create_ppt(title, points, filename="output.pptx"):
    prs = Presentation()

    # Brand Colors
    PRIMARY_PURPLE = RGBColor(94, 42, 132)
    SECONDARY_TEAL = RGBColor(0, 185, 163)
    TEXT_DARK = RGBColor(40, 40, 40)
    BG_LIGHT = RGBColor(244, 244, 244)

    title = clean_title_text(title)

    # Title Slide
    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_PURPLE

    # Title TextBox
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Content Slides
    for idx, item in enumerate(points, start=1):
        key_point = clean_title_text(item.get("title", ""))
        description = item.get("description", "")

        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Alternate background
        bg_color = BG_LIGHT if idx % 2 == 0 else RGBColor(255, 255, 255)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Title
        left, top, width, height = Inches(0.8), Inches(0.5), Inches(8), Inches(1.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.add_paragraph()
        p.text = key_point
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_PURPLE
        p.alignment = PP_ALIGN.LEFT

        # Accent underline
        shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.6), Inches(3), Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SECONDARY_TEAL
        shape.line.fill.background()

        # Description bullets
        if description:
            left, top, width, height = Inches(1), Inches(2.2), Inches(5), Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            for line in description.split("\n"):
                if line.strip():
                    bullet = tf.add_paragraph()
                    bullet.text = line.strip()
                    bullet.font.size = Pt(22)
                    bullet.font.color.rgb = TEXT_DARK
                    bullet.level = 0

        # Footer watermark
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(8), Inches(0.3))
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = "Generated with AI"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.RIGHT

    prs.save(filename)
    return filename

# ---------------- STREAMLIT UI ----------------
st.set_page_config(page_title="AI Productivity Suite", layout="wide")
st.title("AI Productivity Suite")

# ---------------- STATE ----------------
defaults = {
    "messages": [],            # general chat
    "outline_chat": None,      # ppt outline
    "generated_files": [],     # past generated files
    "summary_text": None,      # uploaded doc summary
    "summary_title": None,     # uploaded doc title
    "doc_chat_history": [],    # chat with doc
}
for key, val in defaults.items():
    if key not in st.session_state:
        st.session_state[key] = val

# ---------------- DISPLAY PAST CHAT ----------------
for role, content in st.session_state.messages:
    with st.chat_message(role):
        st.markdown(content)

for role, content in st.session_state.doc_chat_history:
    with st.chat_message(role):
        st.markdown(content)

# ---------------- FILE UPLOAD SECTION ----------------
uploaded_file = st.file_uploader("📂 Upload a document", type=["pdf", "docx", "txt"])

if uploaded_file is not None:
    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        tmp.write(uploaded_file.getvalue())
        tmp_path = tmp.name
    try:
        text = extract_text(tmp_path, uploaded_file.name)
    finally:
        try: os.remove(tmp_path)
        except Exception: pass
    if not text or not text.strip():
        st.error("Unsupported, empty, or unreadable file content.")
    else:
        try:
            summary = summarize_long_text(text)
            title = generate_title(summary) or os.path.splitext(uploaded_file.name)[0]
            st.session_state.summary_text = summary
            st.session_state.summary_title = title
            st.success(f"✅ Document uploaded! Suggested Title: **{st.session_state.summary_title}**. You can now chat with it.")
        except Exception as e:
            st.error(f"Summarization failed: {e}")

# ---------------- ONE CHAT INPUT ----------------
if prompt := st.chat_input("💬 Type a message (general chat or ask about uploaded doc)..."):

    if st.session_state.summary_text:
        # Doc chat mode
        if any(word in prompt.lower() for word in ["ppt", "slides", "presentation"]):
            # User asked for PPT from doc → generate ppt outline
            with st.spinner("Generating PPT outline from uploaded document..."):
                outline = generate_outline_from_desc(st.session_state.summary_text + "\n\n" + prompt)
                outline_data = {"title": st.session_state.summary_title, "slides": outline}
                st.session_state.outline_chat = outline_data
                st.session_state.doc_chat_history.append(("assistant", "✅ Generated PPT outline from document. Preview below."))
        else:
            # Normal doc chat
            st.session_state.doc_chat_history.append(("user", prompt))
            try:
                resp = call_vertex(f"""
You are an assistant answering based only on the provided document.
Document:
{st.session_state.summary_text}

Question:
{prompt}

Answer clearly and concisely using only the document content.
""")
                st.session_state.doc_chat_history.append(("assistant", resp))
            except Exception as e:
                st.session_state.doc_chat_history.append(("assistant", f"⚠️ Backend error: {e}"))

    else:
        # Normal chat / PPT requests
        st.session_state.messages.append(("user", prompt))
        text = prompt.lower()

        try:
            if "ppt" in text or "presentation" in text or "slides" in text:
                with st.spinner("Generating PPT outline..."):
                    outline = generate_outline_from_desc(prompt)
                    st.session_state.outline_chat = {"title": generate_title(prompt), "slides": outline}
                    st.session_state.messages.append(("assistant", "✅ PPT outline generated! Preview below."))
            else:
                bot_reply = call_vertex(prompt)
                st.session_state.messages.append(("assistant", bot_reply))
        except Exception as e:
            st.session_state.messages.append(("assistant", f"⚠️ Backend error: {e}"))

    st.rerun()

# ---------------- OUTLINE PREVIEW + ACTIONS ----------------
if st.session_state.outline_chat:
    outline = st.session_state.outline_chat
    st.subheader(f"📝 Preview Outline: {outline.get('title','Untitled')}")
    for idx, slide in enumerate(outline.get("slides", []), start=1):
        with st.expander(f"Slide {idx}: {slide.get('title', f'Slide {idx}')}", expanded=False):
            st.markdown(slide.get("description", "").replace("\n", "\n\n"))

    # Title edit + feedback area
    new_title = st.text_input("📌 Edit Title", value=outline.get("title", "Untitled"))
    feedback_box = st.text_area("✏️ Feedback for outline (optional):", value="")

    col1, col2 = st.columns(2)

    with col1:
        if st.button("🔄 Apply Feedback"):
            if not feedback_box.strip():
                st.warning("Add feedback text before applying.")
            else:
                with st.spinner("Applying feedback to the outline..."):
                    try:
                        updated = edit_outline_with_feedback(outline, feedback_box)
                        # allow title change
                        updated["title"] = new_title.strip() if new_title.strip() else updated.get("title", "")
                        st.session_state.outline_chat = updated
                        st.success("✅ Outline updated with feedback!")
                        st.experimental_rerun()
                    except Exception as e:
                        st.error(f"❌ Edit failed: {e}")

    with col2:
        if st.button("✅ Generate PPT"):
            with st.spinner("Generating PPT..."):
                try:
                    outline_to_send = copy.deepcopy(outline)
                    outline_to_send["title"] = new_title.strip() if new_title else outline_to_send["title"]

                    filename = f"{re.sub(r'[^A-Za-z0-9_.-]', '_', outline_to_send['title'])}.pptx"
                    ppt_path = create_ppt(outline_to_send["title"], outline_to_send["slides"], filename)
                    with open(ppt_path, "rb") as f:
                        st.download_button(
                            "⬇️ Download PPT",
                            data=f,
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )
                    st.session_state.outline_chat = None
                    st.success("✅ PPT generated and ready to download.")
                except Exception as e:
                    st.error(f"❌ PPT generation error: {e}")
